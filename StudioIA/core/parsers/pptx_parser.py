"""
core/parsers/pptx_parser.py
Estrazione da file PowerPoint (.pptx): testo delle slide, note del relatore,
e immagini incorporate. Per le immagini, il "contesto vicino" è semplicemente
tutto il testo della stessa slide: è un'associazione naturale e affidabile,
molto più semplice che nei documenti Word.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.models import ExtractedImage, ParsedDocument


def extract(file_path: str) -> ParsedDocument:
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return ParsedDocument(source_file=file_path, error=f"PowerPoint illeggibile: {e}")

    full_text_parts = []
    images = []

    for slide_index, slide in enumerate(prs.slides):
        slide_label = f"slide {slide_index + 1}"
        slide_text_parts = []
        picture_shapes = []

        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_text_parts.append(shape.text_frame.text.strip())
            except Exception:
                pass
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_shapes.append(shape)
            except Exception:
                pass

        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_text_parts.append(f"[Note relatore] {notes.strip()}")
        except Exception:
            pass

        slide_text = "\n".join(slide_text_parts)
        if slide_text.strip():
            full_text_parts.append(f"[{slide_label}]\n{slide_text}")

        for pic_idx, shape in enumerate(picture_shapes):
            try:
                image = shape.image
                images.append(
                    ExtractedImage(
                        source_file=file_path,
                        location_hint=f"{slide_label}, immagine {pic_idx + 1}",
                        image_bytes=image.blob,
                        mime_type=image.content_type,
                        nearby_text=slide_text[:1500],
                    )
                )
            except Exception:
                continue

    return ParsedDocument(source_file=file_path, text="\n\n".join(full_text_parts), images=images)
